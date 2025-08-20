# app.py
# Run with: streamlit run app.py
# -----------------------------------------------------------
# NovaMind RAG Chat — Analyze PDFs, DOCX, PPTX, and TXT files
# -----------------------------------------------------------

import os
import io
import re
import math
from typing import List, Dict, Any, Optional, Tuple

import streamlit as st
from dotenv import load_dotenv

# ---------- Optional deps with friendly errors ----------
MISSING_DEPS = []

try:
    from PyPDF2 import PdfReader
except Exception:
    PdfReader = None
    MISSING_DEPS.append("PyPDF2")

try:
    import docx  # python-docx
except Exception:
    docx = None
    MISSING_DEPS.append("python-docx")

try:
    from pptx import Presentation  # python-pptx
except Exception:
    Presentation = None
    MISSING_DEPS.append("python-pptx")

try:
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.metrics.pairwise import cosine_similarity
except Exception:
    TfidfVectorizer = None
    cosine_similarity = None
    MISSING_DEPS.append("scikit-learn")

# OpenRouter (OpenAI SDK)
try:
    from openai import OpenAI
except Exception:
    OpenAI = None
    MISSING_DEPS.append("openai")

# ---------- Setup ----------
load_dotenv()

BOT_NAME = "NovaMind"
DEFAULT_CHAT_MODEL = "qwen/qwen-2.5-7b-instruct"
DEFAULT_EMBED_MODEL = "openai/text-embedding-3-large"  # Will use via OpenRouter if available

def get_api_key() -> Optional[str]:
    # 1) Prefer .env / environment
    key = os.getenv("AI_API_KEY")
    if key:
        return key
    # 2) Try Streamlit secrets ONLY if available
    try:
        return st.secrets["AI_API_KEY"]
    except Exception:
        return None

AI_API_KEY = get_api_key()

st.set_page_config(page_title=f"{BOT_NAME} — RAG Chat", page_icon="📄", layout="wide")
st.markdown(
    """
    <style>
      .block-container { max-width: 1100px; }
      .stChatMessage p { font-size: 1rem; line-height: 1.6; }
      .stChatMessage { border-radius: 12px; }
      .stTextInput > div > div > input { font-size: 1rem; }
      .chunk { background:#f7f7f9; border:1px solid #eee; padding:10px; border-radius:8px; margin-bottom:8px; }
      code.ctx { background:#fff8e1; padding:2px 4px; border-radius:4px; }
    </style>
    """,
    unsafe_allow_html=True,
)

st.title(f"{BOT_NAME} — Document QA")
st.caption("Upload PDFs, DOCX, PPTX, or TXT. Ask questions and get answers grounded in your files.")

# ---------- Dependency checks ----------
if MISSING_DEPS:
    st.warning(
        "The following packages are missing and recommended:\n\n- " +
        "\n- ".join(MISSING_DEPS) +
        "\n\nInstall with:\n\n```bash\npip install streamlit python-dotenv openai PyPDF2 python-docx python-pptx scikit-learn\n```",
        icon="⚠️",
    )

# ---------- API / Client ----------
if AI_API_KEY and OpenAI is not None:
    client = OpenAI(
        base_url="https://openrouter.ai/api/v1",
        api_key=AI_API_KEY,
        default_headers={
            "HTTP-Referer": "http://localhost:8501",
            "X-Title": "NovaMind RAG Chat",
        },
    )
else:
    client = None

# ---------- Sidebar ----------
with st.sidebar:
    st.header("Settings")
    chat_model = st.text_input("Chat model", value=DEFAULT_CHAT_MODEL, help="Any OpenRouter-supported chat model.")
    temperature = st.slider("Creativity (temperature)", 0.0, 1.0, 0.2, 0.1)

    st.subheader("Embeddings (optional)")
    use_embeddings = st.checkbox("Use embeddings for retrieval (faster & smarter)", value=True,
                                 help="If unavailable, the app falls back to TF-IDF.")
    embed_model = st.text_input("Embedding model", value=DEFAULT_EMBED_MODEL,
                                help="OpenAI-style embeddings through OpenRouter, e.g. openai/text-embedding-3-large.")

    chunk_size = st.slider("Chunk size (characters)", 400, 2000, 1200, 100)
    chunk_overlap = st.slider("Chunk overlap (characters)", 0, 400, 200, 50)
    top_k = st.slider("Top-K chunks for context", 2, 10, 5, 1)

    system_prompt_default = (
        f"You are {BOT_NAME}, a helpful assistant. Use ONLY the provided context to answer. "
        f"If the answer is not in the context, say you don't have enough information. "
        f"Be concise, clear, and cite sources like (filename:part)."
    )
    edit_sys = st.checkbox("Edit system prompt", value=False)
    system_prompt = st.text_area("System prompt", value=system_prompt_default, height=120) if edit_sys else system_prompt_default

    st.divider()
    st.caption("Tip: Set AI_API_KEY in `.env` or `.streamlit/secrets.toml`.")

# ---------- Session State ----------
if "index" not in st.session_state:
    st.session_state.index = {
        "chunks": [],          # List[Dict]: {text, source, part_id}
        "vecs": None,          # numpy array or TF-IDF matrix
        "mode": None,          # "emb" or "tfidf"
        "vectorizer": None,    # TF-IDF vectorizer (if used)
    }

if "messages" not in st.session_state:
    st.session_state.messages = [{"role": "system", "content": system_prompt}]
# Keep system message synced
if st.session_state.messages and st.session_state.messages[0]["role"] == "system":
    st.session_state.messages[0]["content"] = system_prompt

# ---------- Utilities ----------
def normalize_ws(s: str) -> str:
    return re.sub(r"\s+", " ", s).strip()

def chunk_text(text: str, max_chars: int = 1200, overlap: int = 200) -> List[str]:
    text = normalize_ws(text)
    if not text:
        return []
    # Prefer splitting on paragraphs/sentences
    paras = re.split(r"(?m)\n{2,}", text)
    chunks = []
    buf = ""
    for p in paras:
        p = p.strip()
        if not p:
            continue
        if len(buf) + len(p) + 1 <= max_chars:
            buf = f"{buf}\n{p}".strip()
        else:
            # Flush buf by size windows
            if buf:
                chunks.extend(_window(buf, max_chars, overlap))
            buf = p
    if buf:
        chunks.extend(_window(buf, max_chars, overlap))
    return [c.strip() for c in chunks if c.strip()]

def _window(s: str, max_chars: int, overlap: int) -> List[str]:
    if len(s) <= max_chars:
        return [s]
    chunks = []
    start = 0
    while start < len(s):
        end = start + max_chars
        chunk = s[start:end]
        chunks.append(chunk)
        if end >= len(s):
            break
        start = end - overlap if overlap > 0 else end
        start = max(start, 0)
    return chunks

# ---------- Extractors ----------
def extract_pdf(file: io.BytesIO) -> List[Tuple[str, str]]:
    """Return list of (part_id, text) per page."""
    if PdfReader is None:
        raise RuntimeError("PyPDF2 not installed.")
    reader = PdfReader(file)
    out = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            txt = page.extract_text() or ""
        except Exception:
            txt = ""
        out.append((f"page {i}", txt))
    return out

def extract_docx(file: io.BytesIO) -> List[Tuple[str, str]]:
    if docx is None:
        raise RuntimeError("python-docx not installed.")
    d = docx.Document(file)
    # Combine paragraphs; use simple 'section X' markers every ~50 paras
    out = []
    buffer, para_count, section_idx = [], 0, 1
    for p in d.paragraphs:
        t = p.text.strip()
        if t:
            buffer.append(t)
            para_count += 1
            if para_count >= 50:
                out.append((f"section {section_idx}", "\n".join(buffer)))
                buffer, para_count, section_idx = [], 0, section_idx + 1
    if buffer:
        out.append((f"section {section_idx}", "\n".join(buffer)))
    # Tables
    tbl_texts = []
    for ti, table in enumerate(d.tables, start=1):
        rows = []
        for row in table.rows:
            cells = [normalize_ws(cell.text) for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            tbl_texts.append("\n".join(rows))
    if tbl_texts:
        out.append((f"tables", "\n\n".join(tbl_texts)))
    return out

def extract_pptx(file: io.BytesIO) -> List[Tuple[str, str]]:
    if Presentation is None:
        raise RuntimeError("python-pptx not installed.")
    prs = Presentation(file)
    out = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = normalize_ws(shape.text)
                if t:
                    texts.append(t)
        # Slide notes
        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = normalize_ws(slide.notes_slide.notes_text_frame.text or "")
            if notes:
                texts.append(f"[Notes] {notes}")
        out.append((f"slide {i}", "\n".join(texts)))
    return out

def extract_txt(file: io.BytesIO) -> List[Tuple[str, str]]:
    data = file.read()
    try:
        txt = data.decode("utf-8")
    except Exception:
        try:
            txt = data.decode("latin-1")
        except Exception:
            txt = data.decode(errors="ignore")
    return [("text", txt)]

def extract_any(uploaded) -> List[Tuple[str, str]]:
    name = uploaded.name.lower()
    data = uploaded.read()
    file = io.BytesIO(data)
    if name.endswith(".pdf"):
        return extract_pdf(file)
    elif name.endswith(".docx"):
        return extract_docx(file)
    elif name.endswith(".pptx"):
        return extract_pptx(file)
    elif name.endswith(".txt"):
        return extract_txt(file)
    else:
        raise RuntimeError("Unsupported file type. Use PDF, DOCX, PPTX, or TXT.")

# ---------- Indexing ----------
def build_index(
    docs: List[Dict[str, Any]],
    use_emb: bool,
    embed_model_name: str,
    chunk_size: int,
    overlap: int,
) -> Dict[str, Any]:
    """
    docs: [{filename, parts: [(part_id, text), ...]}]
    Returns index dict.
    """
    chunks = []
    for d in docs:
        fname = d["filename"]
        for part_id, raw in d["parts"]:
            for i, ch in enumerate(chunk_text(raw, chunk_size, overlap), start=1):
                chunks.append({
                    "text": ch,
                    "source": fname,
                    "part_id": part_id,
                    "chunk_no": i,
                })

    index = {"chunks": chunks, "vecs": None, "mode": None, "vectorizer": None}

    if not chunks:
        return index

    texts = [c["text"] for c in chunks]

    # Try embeddings via OpenRouter; fallback to TF-IDF
    if use_emb and client is not None:
        try:
            # Batch to avoid super long input
            embs = batch_embed(texts, embed_model_name)
            index["vecs"] = embs  # list of lists
            index["mode"] = "emb"
            return index
        except Exception as e:
            st.warning(f"Embedding fallback to TF-IDF due to error: {e}")

    # TF-IDF fallback
    if TfidfVectorizer is None:
        st.error("TF-IDF requires scikit-learn. Please install it or enable embeddings.")
        return index

    vec = TfidfVectorizer(ngram_range=(1, 2), max_features=50_000)
    mat = vec.fit_transform(texts)  # sparse matrix
    index["vecs"] = mat
    index["mode"] = "tfidf"
    index["vectorizer"] = vec
    return index

def batch_embed(texts: List[str], model_name: str, batch_size: int = 64):
    """Embed a list of texts using OpenRouter embeddings (OpenAI-compatible)."""
    vectors = []
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i+batch_size]
        resp = client.embeddings.create(model=model_name, input=batch)
        # OpenAI-compatible: resp.data[i].embedding
        for item in resp.data:
            vectors.append(item.embedding)
    return vectors

def embed_one(text: str, model_name: str) -> List[float]:
    resp = client.embeddings.create(model=model_name, input=[text])
    return resp.data[0].embedding

def retrieve(index: Dict[str, Any], query: str, k: int, emb_model_name: str) -> List[Dict[str, Any]]:
    if not index or not index.get("chunks"):
        return []

    if index["mode"] == "emb":
        try:
            qvec = embed_one(query, emb_model_name)
            # cosine similarity: we can compute manually
            import numpy as np
            A = np.array(index["vecs"], dtype="float32")
            q = np.array(qvec, dtype="float32")
            # Normalize
            A_norm = A / (np.linalg.norm(A, axis=1, keepdims=True) + 1e-8)
            q_norm = q / (np.linalg.norm(q) + 1e-8)
            sims = (A_norm @ q_norm)
            top_idx = sims.argsort()[-k:][::-1]
            return [index["chunks"][i] | {"score": float(sims[i])} for i in top_idx]
        except Exception as e:
            st.warning(f"Embedding retrieval failed, attempting TF-IDF fallback in-memory: {e}")

    # TF-IDF branch
    if index["mode"] == "tfidf":
        if TfidfVectorizer is None or cosine_similarity is None:
            st.error("TF-IDF retrieval requires scikit-learn.")
            return []
        vec = index["vectorizer"]
        mat = index["vecs"]
        q = vec.transform([query])
        sims = cosine_similarity(q, mat)[0]
        top_idx = sims.argsort()[-k:][::-1]
        return [index["chunks"][i] | {"score": float(sims[i])} for i in top_idx]

    return []

def make_context(chunks: List[Dict[str, Any]]) -> str:
    blocks = []
    for c in chunks:
        src = f"{c['source']}:{c['part_id']}#chunk{c['chunk_no']}"
        blocks.append(f"[{src}]\n{c['text']}")
    return "\n\n---\n\n".join(blocks)

def call_llm(messages: List[Dict[str, str]], model: str, temperature: float):
    if client is None:
        raise RuntimeError("Client not initialized (missing AI_API_KEY or openai package).")
    return client.chat.completions.create(
        model=model,
        temperature=temperature,
        stream=True,
        messages=messages,
    )

# ---------- UI: Uploader & Index Builder ----------
st.subheader("1) Upload documents")
uploaded_files = st.file_uploader(
    "Drop files here (PDF, DOCX, PPTX, TXT). You can upload multiple files.",
    type=["pdf", "docx", "pptx", "txt"],
    accept_multiple_files=True
)

colA, colB, colC = st.columns([1,1,1])
with colA:
    build_btn = st.button("Build / Rebuild Index", type="primary")
with colB:
    clear_btn = st.button("Clear Index")
with colC:
    show_chunks = st.checkbox("Preview top chunks after search", value=True)

if clear_btn:
    st.session_state.index = {"chunks": [], "vecs": None, "mode": None, "vectorizer": None}
    st.toast("Index cleared.")

if build_btn:
    if not uploaded_files:
        st.error("Please upload at least one file.")
    else:
        docs = []
        for uf in uploaded_files:
            try:
                parts = extract_any(uf)
                docs.append({"filename": uf.name, "parts": parts})
            except Exception as e:
                st.error(f"Failed to extract {uf.name}: {e}")
        with st.spinner("Building index..."):
            st.session_state.index = build_index(
                docs=docs,
                use_emb=use_embeddings,
                embed_model_name=embed_model,
                chunk_size=chunk_size,
                overlap=chunk_overlap,
            )
        mode = st.session_state.index.get("mode") or "none"
        st.success(f"Index built. Mode: {mode.upper()}. Chunks: {len(st.session_state.index.get('chunks', []))}")

# ---------- UI: Ask Questions ----------
st.subheader("2) Ask questions about your documents")
question = st.chat_input("Ask anything about the uploaded files…")

# Render prior history (skip system)
for m in st.session_state.messages[1:]:
    with st.chat_message(m["role"]):
        st.markdown(m["content"])

if question:
    # Append user message
    st.session_state.messages.append({"role": "user", "content": question})
    with st.chat_message("user"):
        st.markdown(question)

    # Retrieve
    with st.chat_message("assistant"):
        index = st.session_state.index
        if not index or not index.get("chunks"):
            st.error("No index built. Upload files and press 'Build / Rebuild Index'.")
        else:
            rel = retrieve(index, question, top_k, embed_model)
            if show_chunks and rel:
                with st.expander("Context used (top chunks)"):
                    for r in rel:
                        st.markdown(
                            f"<div class='chunk'><strong>{r['source']}</strong> — "
                            f"<code class='ctx'>{r['part_id']}#chunk{r['chunk_no']}</code><br/>"
                            f"<small>score: {r.get('score', 0):.4f}</small><br/><br/>{r['text']}</div>",
                            unsafe_allow_html=True
                        )
            context = make_context(rel)
            sys_msg = {"role": "system", "content": system_prompt}
            user_msg = {
                "role": "user",
                "content": (
                    "Answer the user's question using ONLY the context below. "
                    "If the answer isn't in the context, say so.\n\n"
                    f"### Context\n{context}\n\n"
                    f"### Question\n{question}"
                )
            }

            # Stream reply
            stream_area = st.empty()
            full_reply = ""
            try:
                stream = call_llm([sys_msg, user_msg], chat_model, temperature)
                for chunk in stream:
                    delta = ""
                    try:
                        delta = chunk.choices[0].delta.content or ""
                    except Exception:
                        msg = getattr(chunk.choices[0], "message", None)
                        delta = getattr(msg, "content", "") if msg else ""
                    if delta:
                        full_reply += delta
                        stream_area.markdown(full_reply)
            except Exception as e:
                full_reply = f"⚠️ LLM error: {e}"
                stream_area.markdown(full_reply)

            # Save assistant message
            st.session_state.messages.append({"role": "assistant", "content": full_reply})

# ---------- Footer ----------
st.divider()
left, right = st.columns([1, 2])
with left:
    if st.button("Clear chat"):
        st.session_state.messages = [{"role": "system", "content": system_prompt}]
        st.toast("Chat cleared.")
with right:
    st.caption("Powered by your local index + OpenRouter via OpenAI SDK")
